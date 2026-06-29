import io
import json
import logging
from typing import List, Dict, Any
from app.config import settings
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger("app.presentation")

class PresentationService:
    def __init__(self):
        self.api_key = settings.GEMINI_API_KEY

    async def generate_outline(self, topic: str, outline: str) -> List[Dict[str, Any]]:
        """
        Ask Gemini to build a clean JSON slide outline structure containing slide titles and bullet content.
        Falls back to a structured mock outline if API key is not configured or calls fail.
        """
        system_instruction = (
            "You are a presentation outline designer. Build structured slide plans.\n"
            "You MUST return output strictly as a valid JSON array of slide objects. Do not include markdown code block syntax (like ```json).\n"
            "Format structure:\n"
            "[\n"
            "  {\n"
            "    \"title\": \"Slide Title\",\n"
            "    \"subtitle\": \"Optional Subtitle\",\n"
            "    \"bullets\": [\"Bullet point 1\", \"Bullet point 2\"]\n"
            "  }\n"
            "]"
        )
        
        prompt = (
            f"Generate a presentation structure for topic: {topic}\n"
            f"Guidelines to cover:\n{outline}\n"
            "Provide exactly 4 slides (Title Slide, Slide 2, Slide 3, Closing slide)."
        )

        mock_fallback = [
            {
                "title": f"Intro: {topic}",
                "subtitle": "AI Business Automation Studio Outline",
                "bullets": ["Executive summary analysis", "Core operational targets roadmap"]
            },
            {
                "title": "Key Objectives & Strategy",
                "subtitle": "Defining the core pillars of operations",
                "bullets": ["Integrated microservice architectures", "Vector indexing semantic searches", "Cost-effective token routing"]
            },
            {
                "title": "Implementation Action Plan",
                "subtitle": "Step-by-step milestones configuration",
                "bullets": ["Upload parsed texts caching ingestion", "Interactive RAG search chats panels", "PowerPoint slides auto compilation"]
            },
            {
                "title": "Conclusion & Review",
                "subtitle": "Milestones summary and next iterations",
                "bullets": ["Execute test checks validating endpoints", "Deploy cloud configurations compose containers"]
            }
        ]

        if not self.api_key or self.api_key.startswith("your-"):
            logger.warning("Gemini key not configured. Yielding mock outline.")
            return mock_fallback

        try:
            from langchain_google_genai import ChatGoogleGenerativeAI
            model = ChatGoogleGenerativeAI(
                model="gemini-1.5-flash",
                google_api_key=self.api_key,
                temperature=0.7
            )
            response = await model.ainvoke([
                ("system", system_instruction),
                ("user", prompt)
            ])
            
            clean_text = response.content.strip()
            # Clean possible markdown block wraps
            if clean_text.startswith("```"):
                clean_text = clean_text.split("```")[1]
                if clean_text.startswith("json"):
                    clean_text = clean_text[4:]
            
            return json.loads(clean_text.strip())
        except Exception as e:
            logger.error(f"Failed to generate slide outline from Gemini: {e}")
            return mock_fallback

    def compile_presentation(self, slides_data: List[Dict[str, Any]]) -> io.BytesIO:
        """
        Compiles the slide outline list into a binary PowerPoint (.pptx) file.
        Utilizes dark-theme styling parameters (dark charcoal slides, neon purple accents, white text).
        """
        prs = Presentation()
        
        # Color definitions
        DARK_CHARCOAL = RGBColor(18, 18, 18)
        NEON_PURPLE = RGBColor(168, 85, 247)
        WHITE = RGBColor(240, 240, 240)
        MUTED_GRAY = RGBColor(156, 163, 175)

        for i, slide_item in enumerate(slides_data):
            # Layout: 0 = Title Slide, 1 = Bullet Slide
            layout_idx = 0 if i == 0 else 1
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            # Apply dark background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = DARK_CHARCOAL

            title_text = slide_item.get("title", "Untitled Slide")
            subtitle_text = slide_item.get("subtitle", "")
            bullets = slide_item.get("bullets", [])

            if layout_idx == 0:
                # Title Slide configurations
                title_shape = slide.shapes.title
                subtitle_shape = slide.placeholders[1]

                if title_shape:
                    title_shape.text = title_text
                    title_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
                    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
                    title_shape.text_frame.paragraphs[0].font.bold = True

                if subtitle_shape:
                    subtitle_shape.text = subtitle_text
                    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = NEON_PURPLE
                    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)
            else:
                # Bullet Slide configurations
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]

                if title_shape:
                    title_shape.text = title_text
                    title_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
                    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
                    title_shape.text_frame.paragraphs[0].font.bold = True

                if body_shape:
                    tf = body_shape.text_frame
                    # Subtitle paragraph at top of body box
                    if subtitle_text:
                        p_sub = tf.paragraphs[0]
                        p_sub.text = subtitle_text
                        p_sub.font.color.rgb = NEON_PURPLE
                        p_sub.font.size = Pt(16)
                        p_sub.font.italic = True
                        
                        # Add bullet points below
                        for bullet in bullets:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                            p.font.color.rgb = WHITE
                            p.font.size = Pt(14)
                    else:
                        # Direct bullets
                        for idx, bullet in enumerate(bullets):
                            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                            p.font.color.rgb = WHITE
                            p.font.size = Pt(14)

        # Output to bytes stream buffer
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        return stream

presentation_service = PresentationService()
