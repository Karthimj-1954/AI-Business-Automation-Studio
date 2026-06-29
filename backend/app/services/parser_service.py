import io
import csv
import logging
from typing import Dict, Any

logger = logging.getLogger("app.parser")

class ParserService:
    def parse_txt(self, file_bytes: bytes) -> str:
        """Parse raw text files."""
        return file_bytes.decode("utf-8", errors="ignore")

    def parse_csv(self, file_bytes: bytes) -> str:
        """Parse CSV files into a structured markdown table representation for better LLM retrieval."""
        text_stream = io.StringIO(file_bytes.decode("utf-8", errors="ignore"))
        reader = csv.reader(text_stream)
        
        rows = list(reader)
        if not rows:
            return ""
            
        markdown_lines = []
        # Header row
        headers = rows[0]
        markdown_lines.append("| " + " | ".join(headers) + " |")
        markdown_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        
        # Data rows
        for row in rows[1:]:
            markdown_lines.append("| " + " | ".join(row) + " |")
            
        return "\n".join(markdown_lines)

    def parse_docx(self, file_bytes: bytes) -> str:
        """Parse DOCX files using python-docx."""
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return "\n".join(full_text)
        except ImportError:
            logger.warning("python-docx is not installed. Returning empty parse.")
            return "[Error: python-docx library missing]"
        except Exception as e:
            logger.error(f"Failed to parse DOCX: {e}")
            raise Exception(f"Failed to parse DOCX: {str(e)}")

    def parse_pptx(self, file_bytes: bytes) -> str:
        """Parse PPTX files using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            full_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {i+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                full_text.append("\n".join(slide_text))
            return "\n\n".join(full_text)
        except ImportError:
            logger.warning("python-pptx is not installed. Returning empty parse.")
            return "[Error: python-pptx library missing]"
        except Exception as e:
            logger.error(f"Failed to parse PPTX: {e}")
            raise Exception(f"Failed to parse PPTX: {str(e)}")

    def parse_image(self, file_bytes: bytes) -> str:
        """Run OCR on image file bytes using Tesseract OCR."""
        try:
            from PIL import Image
            import pytesseract
            image = Image.open(io.BytesIO(file_bytes))
            return pytesseract.image_to_string(image)
        except ImportError:
            logger.warning("PIL or pytesseract not installed. OCR failed.")
            return "[Error: OCR libraries (pytesseract/PIL) missing]"
        except Exception as e:
            logger.error(f"Tesseract OCR failed to run: {e}. Ensure Tesseract is installed on host.")
            return f"[OCR Error: Tesseract binary not found or failed: {str(e)}]"

    def parse_pdf(self, file_bytes: bytes) -> str:
        """
        Parse PDF files. Extracts digital text using pypdf.
        Falls back to Tesseract OCR via pdf2image if digital text is empty (scanned PDF).
        """
        extracted_text = ""
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            text_runs = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_runs.append(page_text)
            extracted_text = "\n\n".join(text_runs)
        except ImportError:
            logger.warning("pypdf is not installed. Falling back to OCR.")
        except Exception as e:
            logger.error(f"Failed digital PDF extraction: {e}")

        # If text is empty, it is likely a scanned PDF. Fallback to image-based OCR
        if not extracted_text.strip():
            logger.info("No digital text extracted. Running scanned PDF OCR pipeline.")
            try:
                from pdf2image import convert_from_bytes
                from PIL import Image
                import pytesseract
                
                # Convert PDF pages to images
                images = convert_from_bytes(file_bytes)
                ocr_runs = []
                for i, img in enumerate(images):
                    page_text = pytesseract.image_to_string(img)
                    ocr_runs.append(f"--- Scanned Page {i+1} ---\n{page_text}")
                extracted_text = "\n\n".join(ocr_runs)
            except Exception as e:
                logger.error(f"Failed scanned PDF OCR conversion: {e}")
                extracted_text = f"[PDF Parse Error: No text could be extracted and OCR fallback failed: {str(e)}]"

        return extracted_text

    def parse_file(self, file_bytes: bytes, file_name: str) -> str:
        """
        Route file binary to correct format parser based on file extension name.
        """
        ext = file_name.split(".")[-1].lower()
        if ext == "txt":
            return self.parse_txt(file_bytes)
        elif ext == "csv":
            return self.parse_csv(file_bytes)
        elif ext in ["doc", "docx"]:
            return self.parse_docx(file_bytes)
        elif ext in ["ppt", "pptx"]:
            return self.parse_pptx(file_bytes)
        elif ext == "pdf":
            return self.parse_pdf(file_bytes)
        elif ext in ["png", "jpg", "jpeg"]:
            return self.parse_image(file_bytes)
        else:
            raise ValueError(f"Unsupported file format extension: .{ext}")

parser_service = ParserService()
